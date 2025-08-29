"""
Exportador de Apresentações - PowerPoint e PDF
Versão: 1.0
Data: 2025-08-29
"""

import os
import json
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import weasyprint
    WEASYPRINT_AVAILABLE = True
except ImportError:
    WEASYPRINT_AVAILABLE = False


class PresentationExporter:
    """Exportador para gerar apresentações PowerPoint e PDF"""
    
    def __init__(self, reports_dir: str = "reports"):
        self.reports_dir = Path(reports_dir)
        self.export_dir = Path("exports")
        self.export_dir.mkdir(exist_ok=True)
        
    def export_to_powerpoint(self, general_report_path: str = None) -> str:
        """Exporta relatório geral para PowerPoint"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx não está instalado. Execute: pip install python-pptx")
            
        if general_report_path is None:
            general_report_path = self.reports_dir / "general_analysis.md"
            
        # Carregar dados do relatório geral
        data = self._load_general_report_data(general_report_path)
        
        # Criar apresentação
        prs = Presentation()
        
        # Slide de título
        self._add_title_slide(prs, data)
        
        # Slide de resumo executivo
        self._add_executive_summary_slide(prs, data)
        
        # Slide de estatísticas gerais
        self._add_statistics_slide(prs, data)
        
        # Slide de distribuição por tipo
        self._add_project_types_slide(prs, data)
        
        # Slides de projetos críticos
        self._add_critical_projects_slides(prs, data)
        
        # Slide de próximos passos
        self._add_next_steps_slide(prs)
        
        # Salvar apresentação
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = self.export_dir / f"cnpj_analysis_presentation_{timestamp}.pptx"
        prs.save(str(output_file))
        
        return str(output_file)
        
    def export_to_pdf(self, general_report_path: str = None) -> str:
        """Exporta relatório geral para PDF"""
        if not WEASYPRINT_AVAILABLE:
            raise ImportError("weasyprint não está instalado. Execute: pip install weasyprint")
            
        if general_report_path is None:
            general_report_path = self.reports_dir / "general_analysis.md"
            
        # Ler conteúdo Markdown
        with open(general_report_path, 'r', encoding='utf-8') as f:
            markdown_content = f.read()
            
        # Converter para HTML
        html_content = self._markdown_to_html(markdown_content)
        
        # Gerar PDF
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_file = self.export_dir / f"cnpj_analysis_report_{timestamp}.pdf"
        
        weasyprint.HTML(string=html_content).write_pdf(str(output_file))
        
        return str(output_file)
        
    def _load_general_report_data(self, report_path: Path) -> Dict[str, Any]:
        """Carrega dados do relatório geral"""
        # Procurar pelo JSON correspondente
        json_path = report_path.with_suffix('.json')
        if json_path.exists():
            with open(json_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        
        # Se não encontrar JSON, tentar extrair dados do Markdown
        return self._extract_data_from_markdown(report_path)
        
    def _extract_data_from_markdown(self, report_path: Path) -> Dict[str, Any]:
        """Extrai dados do relatório Markdown"""
        with open(report_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
        data = {
            'total_projects': 0,
            'critical_projects': 0,
            'statistics': {},
            'critical_projects_details': []
        }
        
        # Extrair informações básicas
        lines = content.split('\n')
        for line in lines:
            if '**Total de Projetos**:' in line:
                data['total_projects'] = int(line.split(':')[1].strip())
            elif '**Projetos com Pontos Críticos**:' in line:
                data['critical_projects'] = int(line.split(':')[1].strip())
            elif '**Total de Campos CNPJ**:' in line:
                data['statistics']['total_fields'] = int(line.split(':')[1].strip())
            elif '**Campos Críticos**:' in line:
                data['statistics']['total_critical_fields'] = int(line.split(':')[1].strip())
            elif '**Campos de Alto Impacto**:' in line:
                data['statistics']['total_high_impact_fields'] = int(line.split(':')[1].strip())
                
        return data
        
    def _add_title_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Adiciona slide de título"""
        slide_layout = prs.slide_layouts[0]  # Layout de título
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Análise CNPJ Alfanumérico"
        subtitle.text = f"Relatório Geral - {data.get('total_projects', 0)} Projetos Analisados"
        
        # Estilizar título
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
    def _add_executive_summary_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Adiciona slide de resumo executivo"""
        slide_layout = prs.slide_layouts[1]  # Layout de título e conteúdo
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "📋 Resumo Executivo"
        
        # Criar conteúdo
        text_frame = content.text_frame
        text_frame.clear()
        
        stats = data.get('statistics', {})
        
        p = text_frame.paragraphs[0]
        p.text = f"• Total de Projetos: {data.get('total_projects', 0)}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = f"• Projetos Críticos: {data.get('critical_projects', 0)}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = f"• Total de Campos CNPJ: {stats.get('total_fields', 0):,}"
        p.font.size = Pt(18)
        
        p = text_frame.add_paragraph()
        p.text = f"• Campos Críticos: {stats.get('total_critical_fields', 0)}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(220, 53, 69)  # Vermelho
        
        p = text_frame.add_paragraph()
        p.text = f"• Campos de Alto Impacto: {stats.get('total_high_impact_fields', 0)}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 193, 7)  # Amarelo
        
    def _add_statistics_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Adiciona slide de estatísticas"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "📊 Estatísticas Detalhadas"
        
        text_frame = content.text_frame
        text_frame.clear()
        
        stats = data.get('statistics', {})
        project_types = stats.get('project_types', {})
        
        p = text_frame.paragraphs[0]
        p.text = "Distribuição por Tipo de Projeto:"
        p.font.size = Pt(20)
        p.font.bold = True
        
        for project_type, type_stats in project_types.items():
            p = text_frame.add_paragraph()
            p.text = f"• {project_type}: {type_stats['count']} projetos"
            p.font.size = Pt(16)
            
            if type_stats['critical_fields'] > 0:
                p.text += f" ({type_stats['critical_fields']} críticos)"
                p.font.color.rgb = RGBColor(220, 53, 69)
                
    def _add_project_types_slide(self, prs: Presentation, data: Dict[str, Any]):
        """Adiciona slide de tipos de projeto"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "🏗️ Tipos de Projeto Analisados"
        
        text_frame = content.text_frame
        text_frame.clear()
        
        project_types = {
            'PHP': ['php_symfony', 'php_hyperf', 'php_laravel'],
            'UI/Frontend': ['ui_react', 'ui_vue', 'ui_angular', 'ui_react_native'],
            'Backend/BFF': ['nest_bff'],
            'ETL': ['etl_pentaho', 'etl_sql', 'etl_python']
        }
        
        for category, types in project_types.items():
            p = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 1 else text_frame.add_paragraph()
            p.text = f"• {category}:"
            p.font.size = Pt(18)
            p.font.bold = True
            
            for project_type in types:
                p = text_frame.add_paragraph()
                p.text = f"  - {project_type}"
                p.font.size = Pt(14)
                p.level = 1
                
    def _add_critical_projects_slides(self, prs: Presentation, data: Dict[str, Any]):
        """Adiciona slides de projetos críticos"""
        critical_projects = data.get('critical_projects_details', [])
        
        if not critical_projects:
            # Slide informando que não há projetos críticos
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "✅ Status dos Projetos"
            
            text_frame = content.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            p.text = "🎉 Nenhum projeto com campos críticos encontrado!"
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(40, 167, 69)  # Verde
            
            p = text_frame.add_paragraph()
            p.text = "Todos os projetos estão em conformidade ou precisam apenas de ajustes menores."
            p.font.size = Pt(18)
            
            return
            
        # Adicionar slide de introdução aos projetos críticos
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "🚨 Projetos que Precisam Atenção"
        
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = f"Encontrados {len(critical_projects)} projetos com campos críticos:"
        p.font.size = Pt(20)
        p.font.bold = True
        
        for i, project in enumerate(critical_projects[:5]):  # Limitar a 5 projetos
            p = text_frame.add_paragraph()
            p.text = f"{i+1}. {project['name']} ({project['type']})"
            p.font.size = Pt(16)
            
            if project['critical_fields'] > 0:
                p.text += f" - {project['critical_fields']} campos críticos"
                p.font.color.rgb = RGBColor(220, 53, 69)
                
        # Adicionar slides individuais para projetos críticos
        for project in critical_projects[:3]:  # Limitar a 3 slides detalhados
            self._add_project_detail_slide(prs, project)
            
    def _add_project_detail_slide(self, prs: Presentation, project: Dict[str, Any]):
        """Adiciona slide detalhado de um projeto"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"📁 {project['name']}"
        
        text_frame = content.text_frame
        text_frame.clear()
        
        p = text_frame.paragraphs[0]
        p.text = f"Tipo: {project['type']}"
        p.font.size = Pt(18)
        p.font.bold = True
        
        p = text_frame.add_paragraph()
        p.text = f"Total de Campos: {project['total_fields']}"
        p.font.size = Pt(16)
        
        p = text_frame.add_paragraph()
        p.text = f"Campos Críticos: {project['critical_fields']}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(220, 53, 69)
        
        p = text_frame.add_paragraph()
        p.text = f"Campos de Alto Impacto: {project['high_impact_fields']}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 193, 7)
        
        # Adicionar detalhes dos campos críticos
        if project.get('critical_details'):
            p = text_frame.add_paragraph()
            p.text = "Campos Críticos Encontrados:"
            p.font.size = Pt(14)
            p.font.bold = True
            
            for detail in project['critical_details'][:3]:  # Limitar a 3 campos
                p = text_frame.add_paragraph()
                p.text = f"• {detail['field_name']} - {detail['action_needed']}"
                p.font.size = Pt(12)
                p.level = 1
                
    def _add_next_steps_slide(self, prs: Presentation):
        """Adiciona slide de próximos passos"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "🚀 Próximos Passos"
        
        text_frame = content.text_frame
        text_frame.clear()
        
        steps = [
            "1. Priorizar projetos com campos críticos",
            "2. Revisar campos de alto impacto",
            "3. Planejar migração de dados",
            "4. Atualizar validações frontend",
            "5. Implementar testes de regressão",
            "6. Treinar equipe sobre novo formato"
        ]
        
        for step in steps:
            p = text_frame.paragraphs[0] if len(text_frame.paragraphs) == 1 else text_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(18)
            
    def _markdown_to_html(self, markdown_content: str) -> str:
        """Converte Markdown para HTML"""
        # Conversão simples de Markdown para HTML
        html = markdown_content
        
        # Títulos
        html = html.replace('# ', '<h1>').replace('\n# ', '</h1>\n<h1>')
        html = html.replace('## ', '<h2>').replace('\n## ', '</h2>\n<h2>')
        html = html.replace('### ', '<h3>').replace('\n### ', '</h3>\n<h3>')
        
        # Listas
        html = html.replace('- ', '<li>').replace('\n- ', '</li>\n<li>')
        
        # Negrito
        html = html.replace('**', '<strong>').replace('**', '</strong>')
        
        # Quebras de linha
        html = html.replace('\n\n', '</p>\n<p>')
        
        # Estrutura HTML completa
        html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="UTF-8">
            <title>Relatório CNPJ Alfanumérico</title>
            <style>
                body {{ 
                    font-family: Arial, sans-serif; 
                    margin: 40px; 
                    font-size: 12px;
                    line-height: 1.4;
                }}
                h1 {{ 
                    color: #003366; 
                    font-size: 24px;
                    margin-top: 20px;
                    margin-bottom: 10px;
                }}
                h2 {{ 
                    color: #0066cc; 
                    font-size: 20px;
                    margin-top: 15px;
                    margin-bottom: 8px;
                }}
                h3 {{ 
                    color: #0099ff; 
                    font-size: 16px;
                    margin-top: 12px;
                    margin-bottom: 6px;
                }}
                strong {{ 
                    color: #dc3545; 
                    font-weight: bold;
                }}
                li {{ 
                    margin: 3px 0; 
                    font-size: 11px;
                }}
                p {{
                    margin: 8px 0;
                    font-size: 11px;
                }}
                table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 10px 0;
                    font-size: 10px;
                }}
                th, td {{
                    border: 1px solid #ddd;
                    padding: 4px;
                    text-align: left;
                }}
                th {{
                    background-color: #f2f2f2;
                    font-weight: bold;
                }}
            </style>
        </head>
        <body>
            <p>{html}</p>
        </body>
        </html>
        """
        
        return html
